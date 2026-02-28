import re
import os
import pdfplumber
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# ======================
# MASA TETAP BM
# ======================
masa_tetap = {
    0: "1.00–1.30 petang",
    1: "12.00–1.00",
    2: "9.10–10.10",
    3: "7.40–8.40",
    4: "7.40–8.40"
}

# ======================
# TARIKH MENGIKUT MINGGU
# ======================
tarikh_mingguan = {
    "05": "9–13/2/2026",
    "06": "16–20/2/2026",
    "07": "23–27/2/2026",
    "08": "2–6/3/2026",
    "09": "9–13/3/2026",
    "10": "16–20/3/2026",
    "11": "23–27/3/2026",
    "12": "6–10/4/2026",
    "13": "13–17/4/2026",
    "14": "20–24/4/2026",
    "15": "27/4–1/5/2026",
    "16": "4–8/5/2026",
    "17": "11–15/5/2026",
    "18": "18–22/5/2026",
    "19": "8–12/6/2026",
    "20": "15–19/6/2026",
    "21": "22–26/6/2026",
    "22": "29/6–3/7/2026",
    "23": "6–10/7/2026",
    "24": "13–17/7/2026",
    "25": "20–24/7/2026",
    "26": "27–31/7/2026",
    "27": "3–7/8/2026",
    "28": "10–14/8/2026",
    "29": "17–21/8/2026",
    "30": "24–28/8/2026",
    "31": "31/8–4/9/2026",
    "32": "7–11/9/2026",
    "33": "14–18/9/2026",
    "34": "21–25/9/2026",
    "35": "28/9–2/10/2026",
    "36": "5–9/10/2026",
    "37": "12–16/10/2026",
    "38": "19–23/10/2026"
}

hari_nama = ["Isnin", "Selasa", "Rabu", "Khamis", "Jumaat"]
FONT_SIZE = 11


def extract(text, start, end=None):
    if end:
        pattern = start + r"(.*?)" + end
    else:
        pattern = start + r"(.*)"
    m = re.search(pattern, text, re.DOTALL)
    return m.group(1).strip() if m else ""


# ======================
# FUNCTION UNTUK JANA RPH
# ======================
def generate_rph(pdf_path, kelas, tarikh, hari_dipilih):

    minggu = os.path.basename(pdf_path).replace("M", "").replace(".pdf", "")
    tarikh = tarikh_mingguan.get(minggu, "")

    if not os.path.exists(pdf_path):
        print("❌ File tidak dijumpai:", pdf_path)
        return None

    with pdfplumber.open(pdf_path) as pdf:

        prs = Presentation("templates/template_rph_bm.pptx")
        slide_pointer = 0

        for pdf_index in hari_dipilih:

            if pdf_index >= len(pdf.pages):
                continue

            masa = masa_tetap.get(pdf_index, "")
            page_text = pdf.pages[pdf_index].extract_text()

            tema = extract(page_text, "TEMA", "TAJUK")
            tajuk = extract(page_text, "TAJUK", "STANDARD KANDUNGAN")
            sk = extract(page_text, "STANDARD KANDUNGAN", "STANDARD PEMBELAJARAN")

            sp_full = extract(page_text, "STANDARD PEMBELAJARAN", "AKTIVITI PERMULAAN")
            sp = sp_full.split("\n")[0] if sp_full else ""

            objektif = extract(page_text, "Pada akhir pengajaran", "AKTIVITI PERMULAAN")

            aktiviti_raw = extract(page_text, "AKTIVITI UTAMA", "AKTIVITI PENUTUP")
            aktiviti_lines = []

            for line in aktiviti_raw.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if "(AKTIVITI)" in line:
                    continue
                if "***Murid" in line:
                    continue
                aktiviti_lines.append(line)

            aktiviti = "\n".join(aktiviti_lines)

            refleksi_full = extract(page_text, "REFLEKSI")
            refleksi = refleksi_full.split("\n")[0] if refleksi_full else ""

            emk = extract(page_text, "EMK:", "PENILAIAN")
            peta = extract(page_text, "PETA", "KB")

            pak21_parts = []
            if emk:
                pak21_parts.append(emk.strip())
            if peta:
                pak21_parts.append(peta.strip())

            pak21 = "\n".join(pak21_parts)

            catatan_section = extract(page_text, "CATATAN", "KAEDAH")
            kehadiran_match = re.search(r"\d+/\d+\s*orang", catatan_section)
            kehadiran = kehadiran_match.group(0) if kehadiran_match else ""

            if slide_pointer >= len(prs.slides):
                break

            slide = prs.slides[slide_pointer]
            slide_pointer += 1

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table

                    for row in table.rows:
                        for cell in row.cells:

                            text = cell.text_frame.text

                            text = text.replace("{{KELAS}}", kelas)
                            text = text.replace("{{MASA}}", masa)
                            text = text.replace("{{TARIKH}}", tarikh)
                            text = text.replace("{{HARI}}", hari_nama[pdf_index])
                            text = text.replace("{{TEMA}}", tema)
                            text = text.replace("{{TAJUK}}", tajuk)
                            text = text.replace("{{SK}}", sk)
                            text = text.replace("{{SP}}", sp)
                            text = text.replace("{{OBJEKTIF}}", objektif)
                            text = text.replace("{{AKTIVITI}}", aktiviti)
                            text = text.replace("{{REFLEKSI}}", refleksi)
                            text = text.replace("{{PAK21}}", pak21)
                            text = text.replace("{{KEHADIRAN}}", kehadiran)

                            cell.text_frame.text = text

                            tf = cell.text_frame
                            tf.word_wrap = True

                            for paragraph in tf.paragraphs:
                                paragraph.space_before = 0
                                paragraph.space_after = 0
                                for run in paragraph.runs:
                                    run.font.size = Pt(FONT_SIZE)

                            if "BAHASA MELAYU" in cell.text:
                                for paragraph in tf.paragraphs:
                                    paragraph.alignment = PP_ALIGN.CENTER

        output_file = "RPH_BM_FINAL_OUTPUT.pptx"
        prs.save(output_file)

    return output_file


# ======================
# TEST MANUAL (SERVER GUNA INI)
# ======================
if __name__ == "__main__":
    import sys
    import json

    minggu = sys.argv[1]
    tarikh = sys.argv[2]
    kelas = sys.argv[3]
    hari_dipilih = json.loads(sys.argv[4])

    pdf_path = f"bm/M{minggu}.pdf"

    output = generate_rph(
        pdf_path,
        kelas,
        tarikh,
        hari_dipilih
    )

    if output:
        print("SIAP")
